"""
Assemble final slides into a PowerPoint presentation.

Usage:
    python assemble.py

Output:
    Claude_Code_Tips_for_Engineers.pptx
"""

from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

SLIDES_DIR = Path(__file__).parent / "output" / "final"
OUTPUT = Path(__file__).parent / "Claude_Code_Tips_for_Engineers.pptx"


def main():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Blank slide layout
    blank_layout = prs.slide_layouts[6]

    # Add all 12 slides
    for i in range(1, 13):
        slide_path = SLIDES_DIR / f"slide_{i:02d}.png"
        if slide_path.exists():
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(slide_path),
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )
            print(f"Added slide {i}")
        else:
            print(f"Missing: slide_{i:02d}.png")

    prs.save(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Size: {OUTPUT.stat().st_size / 1024 / 1024:.1f} MB")
    print("\nUpload to Google Slides:")
    print("1. Go to slides.google.com")
    print("2. File → Import slides → Upload")
    print("3. Select all slides → Import")


if __name__ == "__main__":
    main()
